"""
Prepare le corpus documentaire a partir de fichiers sources bruts, quel que
soit leur format (Word .doc/.docx, PDF, Excel, texte/Markdown), deposes dans
data/source_documents/. Convertit chacun en texte brut dans data/docs/, pret
a etre decoupe en chunks et indexe par ingest.py.

Workflow pour l'equipe OPO : pour ajouter un nouveau document de reference
(manuel, fiche, note, guide de procedure...), il suffit de le deposer dans
data/source_documents/ puis de relancer `python ingest.py` (qui appelle
automatiquement ce module avant de reconstruire l'index) - une seule commande,
aucune connaissance technique du format necessaire.

Formats geres : .doc, .docx, .pdf, .ppt, .pptx, .txt, .md, .xlsx (feuilles non structurees
en dictionnaire de variables - le fichier `Dictionnaire_donnees_OPO_*.xlsx`
continue d'etre traite a part par la logique existante, qui produit un
document par table avec un decoupage variable-par-variable, de bien meilleure
qualite que ne le ferait une extraction generique).

Chaque fichier source converti ne l'est qu'une seule fois (verification par
date de modification) : reconvertir un PDF ou un .doc via LibreOffice est lent,
inutile de le refaire a chaque lancement si le fichier source n'a pas change.
"""

import re
import subprocess
import tempfile
from pathlib import Path

SOURCE_DIR = Path(__file__).parent / "data" / "source_documents"
DOCS_DIR = Path(__file__).parent / "data" / "docs"

EXTENSIONS_SUPPORTEES = {".doc", ".docx", ".pdf", ".txt", ".md", ".xlsx", ".ppt", ".pptx"}

# Fichiers a ignorer s'ils se retrouvent deposes ici par erreur (le dictionnaire
# de variables est deja traite specifiquement ailleurs, cf. docstring ci-dessus).
NOMS_IGNORES = {"dictionnaire_donnees_opo_basedemographique"}


def _nettoyer_texte(texte: str) -> str:
    """Normalise les fins de ligne et evite les rafales de plus de 2 lignes
    vides, sans casser les separations de paragraphes dont depend le
    decoupage en chunks (ingest.chunk_file coupe sur les lignes vides)."""
    lignes = [l.rstrip() for l in texte.replace("\r\n", "\n").splitlines()]
    texte = "\n".join(lignes)
    texte = re.sub(r"\n{3,}", "\n\n", texte)
    return texte.strip()


def _convertir_via_soffice(path: Path) -> str:
    """Conversion generique via LibreOffice en mode headless : gere aussi
    bien les .doc (ancien format binaire, aucune librairie Python pure ne le
    lit correctement) que les .docx en repli si python-docx echoue."""
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmp, str(path)],
            check=True, capture_output=True, timeout=180,
        )
        sortie = Path(tmp) / (path.stem + ".txt")
        if not sortie.exists():
            raise RuntimeError(f"LibreOffice n'a pas produit de fichier texte pour {path.name}")
        return sortie.read_text(encoding="utf-8", errors="ignore")


def extraire_doc(path: Path) -> str:
    """Fichier Word ancien format (.doc, Word 97-2003) : necessite LibreOffice
    (soffice) installe sur la machine - c'est le cas sur le serveur de
    deploiement prevu (VPS Linux), a installer sur un poste Windows sans Word
    si besoin (LibreOffice est gratuit)."""
    return _convertir_via_soffice(path)


def extraire_docx(path: Path) -> str:
    """Fichier Word moderne (.docx) : lecture directe via python-docx (rapide,
    pas de dependance externe), avec repli sur LibreOffice si la lecture
    directe echoue (fichier atypique ou legerement corrompu)."""
    try:
        import docx

        document = docx.Document(str(path))
        paragraphes = [p.text for p in document.paragraphs]
        # Les tableaux contiennent souvent l'essentiel de l'info dans les
        # fiches techniques (ex: tableaux de correspondance) - a inclure.
        for table in document.tables:
            for row in table.rows:
                ligne = " | ".join(cell.text for cell in row.cells)
                if ligne.strip():
                    paragraphes.append(ligne)
        texte = "\n\n".join(p for p in paragraphes if p.strip())
        if not texte.strip():
            raise ValueError("aucun texte trouve via python-docx")
        return texte
    except Exception:
        return _convertir_via_soffice(path)


def extraire_pdf(path: Path) -> str:
    """Fichier PDF : extraction via pdftotext (poppler-utils, sans l'option
    -layout pour laisser le texte se re-agencer en paragraphes naturels -
    meilleur pour le decoupage par blocs que le rendu colonne-par-colonne).
    Repli sur pdfplumber (pure Python, pip) si pdftotext n'est pas installe."""
    try:
        resultat = subprocess.run(
            ["pdftotext", str(path), "-"],
            check=True, capture_output=True, timeout=180,
        )
        return resultat.stdout.decode("utf-8", errors="ignore")
    except (FileNotFoundError, subprocess.CalledProcessError):
        import pdfplumber

        morceaux = []
        with pdfplumber.open(str(path)) as pdf:
            for page in pdf.pages:
                morceaux.append(page.extract_text() or "")
        return "\n\n".join(morceaux)


def extraire_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _texte_pptx(chemin_pptx: Path) -> str:
    """Extrait le texte d'un fichier .pptx (deja au format XML moderne) : une
    diapositive = un bloc, texte des zones de texte et des tableaux."""
    from pptx import Presentation

    presentation = Presentation(str(chemin_pptx))
    blocs = []
    for i, slide in enumerate(presentation.slides, start=1):
        morceaux = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                morceaux.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    ligne = " | ".join(cell.text for cell in row.cells)
                    if ligne.strip():
                        morceaux.append(ligne)
        if morceaux:
            blocs.append(f"Diapositive {i} : " + " ".join(morceaux))
    return "\n\n".join(blocs)


def extraire_pptx(path: Path) -> str:
    """Fichier PowerPoint moderne (.pptx) : lecture directe via python-pptx,
    avec repli sur LibreOffice (reconversion en .pptx propre) si la lecture
    directe echoue."""
    try:
        texte = _texte_pptx(path)
        if not texte.strip():
            raise ValueError("aucun texte trouve via python-pptx")
        return texte
    except Exception:
        with tempfile.TemporaryDirectory() as tmp:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pptx", "--outdir", tmp, str(path)],
                check=True, capture_output=True, timeout=180,
            )
            reconverti = Path(tmp) / (path.stem + ".pptx")
            return _texte_pptx(reconverti)


def extraire_ppt(path: Path) -> str:
    """Fichier PowerPoint ancien format (.ppt) : LibreOffice ne sait pas
    l'exporter directement en texte (pas de filtre "txt" pour Impress), donc
    on le reconvertit d'abord en .pptx, puis on relit ce .pptx via
    python-pptx (memes moyens que pour un .pptx deja moderne)."""
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pptx", "--outdir", tmp, str(path)],
            check=True, capture_output=True, timeout=180,
        )
        converti = Path(tmp) / (path.stem + ".pptx")
        if not converti.exists():
            raise RuntimeError(f"LibreOffice n'a pas produit de .pptx pour {path.name}")
        return _texte_pptx(converti)


def extraire_xlsx_generique(path: Path) -> str:
    """Extraction generique d'un classeur Excel qui n'a pas la structure du
    dictionnaire de variables (une feuille = une liste de variables) - utile
    pour un tableau de suivi, une grille de codes, etc. deposee comme document
    de reference plutot que comme table de donnees a analyser."""
    import openpyxl

    classeur = openpyxl.load_workbook(path, data_only=True, read_only=True)
    blocs = []
    for nom_feuille in classeur.sheetnames:
        feuille = classeur[nom_feuille]
        lignes_texte = []
        for ligne in feuille.iter_rows(values_only=True):
            valeurs = [str(v) for v in ligne if v is not None]
            if valeurs:
                lignes_texte.append(" | ".join(valeurs))
        if lignes_texte:
            blocs.append(f"Feuille: {nom_feuille}\n" + "\n".join(lignes_texte))
    return "\n\n".join(blocs)


def convertir_fichier(path: Path) -> str | None:
    ext = path.suffix.lower()
    if ext == ".doc":
        return extraire_doc(path)
    if ext == ".docx":
        return extraire_docx(path)
    if ext == ".pdf":
        return extraire_pdf(path)
    if ext in (".txt", ".md"):
        return extraire_txt(path)
    if ext == ".xlsx":
        return extraire_xlsx_generique(path)
    if ext == ".pptx":
        return extraire_pptx(path)
    if ext == ".ppt":
        return extraire_ppt(path)
    return None


def preparer_corpus() -> list[str]:
    """Convertit tous les fichiers de data/source_documents/ en texte dans
    data/docs/ (prefixe 'source_' pour ne jamais ecraser les documents deja
    geres a la main, comme les fiches par table du dictionnaire ou le schema
    relationnel). Ignore un fichier si sa conversion existe deja et est plus
    recente que le fichier source. Renvoie la liste des fichiers convertis."""
    SOURCE_DIR.mkdir(parents=True, exist_ok=True)
    DOCS_DIR.mkdir(parents=True, exist_ok=True)

    convertis = []
    for path in sorted(SOURCE_DIR.iterdir()):
        if not path.is_file() or path.suffix.lower() not in EXTENSIONS_SUPPORTEES:
            continue
        if path.stem.lower() in NOMS_IGNORES:
            continue

        sortie = DOCS_DIR / f"source_{path.stem}.txt"
        if sortie.exists() and sortie.stat().st_mtime >= path.stat().st_mtime:
            continue

        try:
            texte = convertir_fichier(path)
        except Exception as e:
            print(f"[prepare_corpus] Echec de conversion pour {path.name} : {e}")
            continue

        if not texte or not texte.strip():
            print(f"[prepare_corpus] {path.name} : aucun texte extrait, fichier ignore.")
            continue

        entete = f"Document: {path.stem}\n\n"
        sortie.write_text(entete + _nettoyer_texte(texte), encoding="utf-8")
        convertis.append(path.name)

    return convertis


if __name__ == "__main__":
    fichiers = preparer_corpus()
    if fichiers:
        print(f"{len(fichiers)} document(s) converti(s) : {', '.join(fichiers)}")
    else:
        print("Aucun nouveau document a convertir dans data/source_documents/.")
